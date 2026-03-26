from fastapi import UploadFile, HTTPException
from core.mongo import documents_col, system_logs_col
from core.minio_client import minio_client, settings
from minio.error import S3Error
from utils.text_process import clean_text
from bson import ObjectId
from datetime import datetime
import os
import PyPDF2
import docx
import pptx


# 解析不同格式文件（PDF/DOC/TXT/PPT，设计文档支持格式）
def parse_file(file_path: str, file_type: str) -> str:
    text = ""
    try:
        if file_type == "pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
        elif file_type == "doc" or file_type == "docx":
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_type == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        elif file_type == "ppt" or file_type == "pptx":
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            raise HTTPException(status_code=400, detail=f"不支持的文件类型：{file_type}")
        return clean_text(text)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件解析失败：{str(e)}")


# 文档上传：MinIO存储+MongoDB元数据记录（设计文档流程）
async def upload_document(file: UploadFile, user_id: str):
    # 1. 格式/大小校验（设计文档要求）
    allowed_types = ["pdf", "doc", "docx", "txt", "ppt", "pptx"]
    file_type = file.filename.split(".")[-1].lower()
    if file_type not in allowed_types:
        raise HTTPException(status_code=400, detail=f"仅支持{allowed_types}格式")
    if file.size > 100 * 1024 * 1024:  # 100MB限制
        raise HTTPException(status_code=400, detail="文件大小不超过100MB")

    # 2. 生成唯一文档ID
    doc_id = str(ObjectId())
    file_name = f"{doc_id}.{file_type}"
    file_path = f"/tmp/{file_name}"

    # 3. 临时保存文件
    with open(file_path, "wb") as f:
        f.write(await file.read())

    # 4. 上传至MinIO
    try:
        minio_client.fput_object(
            bucket_name=settings.MINIO_BUCKET,
            object_name=file_name,
            file_path=file_path
        )
    except S3Error as e:
        os.remove(file_path)
        raise HTTPException(status_code=500, detail=f"文件存储失败：{str(e)}")

    # 5. 写入MongoDB元数据（设计文档documents集合）
    doc_meta = {
        "_id": doc_id,
        "name": file.filename,
        "type": file_type,
        "size": file.size,
        "upload_time": datetime.now(),
        "update_time": datetime.now(),
        "upload_user_id": user_id,
        "status": "pending",  # 待处理
        "parse_log": ""
    }
    documents_col.insert_one(doc_meta)

    # 6. 记录系统日志
    system_logs_col.insert_one({
        "module": "document",
        "operation": "upload",
        "user_id": user_id,
        "content": f"用户{user_id}上传文档：{file.filename}，ID：{doc_id}",
        "create_time": datetime.now(),
        "level": "info"
    })

    # 7. 删除临时文件
    os.remove(file_path)

    return {"doc_id": doc_id, "file_name": file.filename, "status": "pending"}


# 文档检索/删除（设计文档要求）
def get_documents(filters: dict = None):
    filters = filters or {}
    return list(documents_col.find(filters, {"_id": 1, "name": 1, "type": 1, "upload_time": 1, "status": 1}))


def delete_document(doc_id: str, user_id: str):
    # 逻辑删除（设计文档要求）
    result = documents_col.update_one(
        {"_id": doc_id},
        {"$set": {"status": "deleted", "update_time": datetime.now()}}
    )
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="文档不存在或已删除")
    # 记录日志
    system_logs_col.insert_one({
        "module": "document",
        "operation": "delete",
        "user_id": user_id,
        "content": f"用户{user_id}删除文档，ID：{doc_id}",
        "create_time": datetime.now(),
        "level": "info"
    })
    return {"status": "success", "doc_id": doc_id}