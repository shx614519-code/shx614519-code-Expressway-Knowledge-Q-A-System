from fastapi import UploadFile, HTTPException
from core.mongo import documents_col, system_logs_col
from core.minio_client import minio_client, settings
from minio.error import S3Error
from utils.text_process import clean_text
from bson import ObjectId
from datetime import datetime
import os
import PyPDF2
import docx
import pptx


# 解析不同格式文件（PDF/DOC/TXT/PPT，设计文档支持格式）
def parse_file(file_path: str, file_type: str) -> str:
    text = ""
    try:
        if file_type == "pdf":
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() or ""
        elif file_type == "doc" or file_type == "docx":
            doc = docx.Document(file_path)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_type == "txt":
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()
        elif file_type == "ppt" or file_type == "pptx":
            prs = pptx.Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        else:
            raise HTTPException(status_code=400, detail=f"不支持的文件类型：{file_type}")
        return clean_text(text)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"文件解析失败：{str(e)}")



# 文档上传：MinIO 存储+MongoDB 元数据记录（设计文档流程）
async def upload_document(file: UploadFile, user_id: str) -> dict:
    from datetime import datetime
    from bson import ObjectId
    import io
    from minio.error import S3Error

    # 生成文档 ID
    doc_id = str(ObjectId())

    # 获取文件扩展名
    file_ext = file.filename.split(".")[-1].lower() if "." in file.filename else "bin"

    try:
        # 读取文件内容
        content = await file.read()

        # 上传到 MinIO
        object_name = f"{doc_id}.{file_ext}"
        # 将字节数据包装成 BytesIO 对象
        data_stream = io.BytesIO(content)

        # 确保 bucket 存在，不存在则创建
        if not minio_client.bucket_exists(settings.MINIO_BUCKET):
            minio_client.make_bucket(settings.MINIO_BUCKET)

        minio_client.put_object(
            settings.MINIO_BUCKET,
            object_name,
            data=data_stream,
            length=len(content)
        )

        # 插入文档记录到 MongoDB
        documents_col.insert_one({
            "_id": doc_id,
            "name": file.filename,
            "type": file_ext,
            "file_path": f"{settings.MINIO_BUCKET}/{object_name}",
            "upload_user_id": user_id,
            "upload_time": datetime.now(),
            "status": "pending",
            "knowledge_count": 0
        })

        # 记录日志
        system_logs_col.insert_one({
            "module": "document",
            "operation": "upload",
            "user_id": user_id,
            "content": f"用户上传文档 {file.filename}，ID: {doc_id}",
            "create_time": datetime.now(),
            "level": "info"
        })

        return {
            "doc_id": doc_id,
            "filename": file.filename,
            "type": file_ext,
            "status": "pending"
        }

    except S3Error as e:
        # MinIO 错误处理
        from fastapi import HTTPException
        raise HTTPException(
            status_code=500,
            detail=f"MinIO 上传失败：{str(e)}"
        )
    except Exception as e:
        # 通用错误处理
        from fastapi import HTTPException
        raise HTTPException(
            status_code=500,
            detail=f"文件上传失败：{str(e)}"
        )


# 文档检索/删除（设计文档要求）
def get_documents(filters: dict = None):
    filters = filters or {}
    return list(documents_col.find(filters, {"_id": 1, "name": 1, "type": 1, "upload_time": 1, "status": 1}))


def delete_document(doc_id: str, user_id: str):
    # 逻辑删除（设计文档要求）
    result = documents_col.update_one(
        {"_id": doc_id},
        {"$set": {"status": "deleted", "update_time": datetime.now()}}
    )
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="文档不存在或已删除")
    # 记录日志
    system_logs_col.insert_one({
        "module": "document",
        "operation": "delete",
        "user_id": user_id,
        "content": f"用户{user_id}删除文档，ID：{doc_id}",
        "create_time": datetime.now(),
        "level": "info"
    })
    return {"status": "success", "doc_id": doc_id}